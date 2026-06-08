#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Builds the VKR defense presentation from vkr_title.pptx + rendered assets."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

AD = "slides_assets"
BLACK = RGBColor(0x15, 0x15, 0x15)
DARK  = RGBColor(0x22, 0x22, 0x22)
GRAY  = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xBB, 0xBB, 0xBB)
FONT  = "Arial"

prs = Presentation("vkr_title.pptx")
SW, SH = prs.slide_width, prs.slide_height          # 13.333 x 7.5 in
SW_IN, SH_IN = SW / 914400, SH / 914400

# pick a blank layout (fewest placeholders)
blank = min(prs.slide_layouts, key=lambda L: len(L.placeholders))

_page = [1]  # title is page 1

def white_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

def add_slide():
    s = prs.slides.add_slide(blank)
    # remove any inherited placeholders
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    white_bg(s)
    _page[0] += 1
    pn = s.shapes.add_textbox(Inches(SW_IN-0.7), Inches(SH_IN-0.55), Inches(0.5), Inches(0.4))
    r = pn.text_frame.paragraphs[0].add_run(); r.text = str(_page[0])
    r.font.size = Pt(14); r.font.name = FONT; r.font.color.rgb = GRAY
    return s

def title(s, text, size=30):
    box = s.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(SW_IN-1.1), Inches(1.05))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.bold = True; r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = BLACK
    return box

def body(s, items, left=0.6, top=1.5, width=None, height=None, size=18, gap=6):
    """items: list of dicts {text, bold, italic, bullet, color, size, space_before}."""
    if width is None: width = SW_IN - 1.2
    if height is None: height = SH_IN - top - 0.4
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(it.get("space_after", gap))
        p.space_before = Pt(it.get("space_before", 0))
        p.line_spacing = it.get("line_spacing", 1.04)
        txt = it["text"]
        if it.get("bullet"):
            txt = "•  " + txt
        r = p.add_run(); r.text = txt
        r.font.size = Pt(it.get("size", size)); r.font.name = FONT
        r.font.bold = it.get("bold", False); r.font.italic = it.get("italic", False)
        r.font.color.rgb = it.get("color", DARK)
        if it.get("bullet"):
            p.level = it.get("level", 0)
    return box

def img_size(path):
    iw, ih = Image.open(path).size
    return iw / ih

def picture(s, name, left, top, max_w, max_h, caption=None):
    path = os.path.join(AD, name + ".png")
    ar = img_size(path)
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        cb = s.shapes.add_textbox(Inches(left), Inches(top + max_h + 0.02), Inches(max_w), Inches(0.4))
        p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(12); r.font.name = FONT; r.font.italic = True; r.font.color.rgb = GRAY
    return (x, y, w, h)

def _set_border(cell, edges=("L","R","T","B"), color="BBBBBB", w=9525):
    tcPr = cell._tc.get_or_add_tcPr()
    tag = {"L":"a:lnL","R":"a:lnR","T":"a:lnT","B":"a:lnB"}
    for e in edges:
        ln = tcPr.find(qn(tag[e]))
        if ln is not None:
            tcPr.remove(ln)
        ln = tcPr.makeelement(qn(tag[e]), {"w": str(w), "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
        fill.append(clr); ln.append(fill)
        # insert in schema order (lnL,lnR,lnT,lnB before fill props) - append works for rendering
        tcPr.append(ln)

def table(s, rows, left, top, width, col_w=None, size=14, header=True,
          first_col_left=True, height_row=0.42):
    nr = len(rows); nc = len(rows[0])
    gframe = s.shapes.add_table(nr, nc, Inches(left), Inches(top),
                                Inches(width), Inches(height_row * nr))
    t = gframe.table
    t.first_row = False; t.horz_banding = False; t.first_col = False
    if col_w:
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            t.columns[i].width = Inches(width * cw / tot)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Inches(0.08); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (ci == 0 and first_col_left) else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = DARK
            bold = (header and ri == 0)
            # allow per-cell bold via **...**
            if isinstance(val, str) and val.startswith("**") and val.endswith("**"):
                r.text = val[2:-2]; bold = True
            r.font.bold = bold
            _set_border(c)
    return gframe

# =====================================================================
# Slide 1 — title (update existing)
# =====================================================================
ts = prs.slides[0]
def set_ph_text(shape, lines, keepfmt=True):
    tf = shape.text_frame
    # use first paragraph's first run formatting as template
    p0 = tf.paragraphs[0]
    base = p0.runs[0] if p0.runs else None
    align = p0.alignment
    sz = base.font.size if base else Pt(20)
    nm = base.font.name if base else FONT
    bold = base.font.bold if base else False
    col = None
    try: col = base.font.color.rgb
    except Exception: col = None
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run(); r.text = ln
        if sz: r.font.size = sz
        r.font.name = nm or FONT
        r.font.bold = bool(bold)
        if col is not None:
            try: r.font.color.rgb = col
            except Exception: pass

for sh in ts.shapes:
    if not sh.has_text_frame: continue
    nm = sh.name
    if nm.startswith("Заголовок"):
        set_ph_text(sh, ["Разработка системы шумоподавления и очистки",
                         "аудиосигналов на основе глубокого обучения"])
    elif nm.startswith("Подзаголовок"):
        set_ph_text(sh, ["Вамбриков Никита Сергеевич", "МФТИ"])
    elif sh.text and "Москва" in sh.text:
        set_ph_text(sh, ["Москва, 2026"])
    elif sh.text and "руководител" in sh.text:
        set_ph_text(sh, ["Научный руководитель: д-р физ.-мат. наук, проф. Цурков Владимир Иванович"])

# =====================================================================
# Slide 2 — Введение (2.1–2.5)
# =====================================================================
s = add_slide(); title(s, "Введение")
body(s, [
    {"text": "Область и востребованность", "bold": True, "size": 18, "space_before": 2},
    {"text": "Улучшение и восстановление речи требуется в синтезе (TTS) и распознавании (ASR) речи, телефонии, слуховых аппаратах, реставрации архивов и автоматическом создании обучающих датасетов студийного качества.", "size": 16, "space_after": 8},
    {"text": "Цель работы", "bold": True, "size": 18},
    {"text": "Построить и исследовать открытую систему восстановления речи на основе архитектуры Miipher-2, заменив недоступный энкодер USM открытыми компонентами.", "size": 16, "space_after": 8},
    {"text": "Существующие подходы", "bold": True, "size": 18},
    {"text": "Классические (спектральное вычитание, фильтр Винера); дискриминативные нейросети; генеративные (GAN, диффузия); восстановление через ресинтез (Miipher, Miipher-2).", "size": 16, "space_after": 8},
    {"text": "Основная трудность", "bold": True, "size": 18},
    {"text": "Лучшая по качеству модель Miipher-2 опирается на закрытый энкодер USM (2 млрд параметров, 12 млн ч обучения) — веса недоступны, воспроизведение невозможно.", "size": 16, "space_after": 8},
    {"text": "Предлагаемый подход", "bold": True, "size": 18},
    {"text": "Заменить USM открытыми SSL-энкодерами (WavLM, HuBERT), исследовать оптимальный слой извлечения признаков и открытые вокодеры (HiFi-GAN, WaveFit).", "size": 16},
], top=1.35)

# =====================================================================
# Slide 3 — Анализ базовых источников
# =====================================================================
s = add_slide(); title(s, "Анализ базовых источников")
body(s, [
    {"text": "Восстановление речи через ресинтез", "bold": True, "size": 18, "space_before": 2},
    {"text": "Koizumi et al. Miipher (WASPAA 2023); Miipher-2 (arXiv 2025). Недостаток: закрытый энкодер USM; у Miipher-1 — текстовое и дикторское кондиционирование.", "bullet": True, "size": 15, "space_after": 8},
    {"text": "Самоконтролируемые энкодеры (SSL)", "bold": True, "size": 18},
    {"text": "HuBERT (2021), WavLM (2022), wav2vec 2.0 (2020). Недостаток: обучены преимущественно на английской речи; разные слои несут разную информацию.", "bullet": True, "size": 15, "space_after": 8},
    {"text": "Нейросетевые вокодеры", "bold": True, "size": 18},
    {"text": "HiFi-GAN (2020), WaveFit (2023), BigVGAN (2023). Недостаток: WaveFit итеративен (медленнее в 5 раз), BigVGAN тяжёлый (112 млн параметров).", "bullet": True, "size": 15, "space_after": 8},
    {"text": "Классические и дискриминативные методы", "bold": True, "size": 18},
    {"text": "Спектральное вычитание, фильтр Винера, DEMUCS. Недостаток: не восстанавливают подавленные гармоники, вносят музыкальный шум.", "bullet": True, "size": 15},
], top=1.35)

# =====================================================================
# Slide 4 — Базовая модель: Miipher-2
# =====================================================================
s = add_slide(); title(s, "Базовая модель: Miipher-2")
picture(s, "fig_pipeline", left=0.7, top=1.45, max_w=12.0, max_h=2.3,
        caption="Рис. 1. Схема восстановления речи в пространстве SSL-признаков")
body(s, [
    {"text": "Очистка выполняется не во временной или спектральной области, а в пространстве признаков SSL-энкодера.", "bullet": True, "size": 16},
    {"text": "Схема: SSL-энкодер → очиститель признаков → нейросетевой вокодер.", "bullet": True, "size": 16},
    {"text": "Miipher-2: замороженный USM + параллельные адаптеры (~20 млн) + вокодер WaveFit; без текста и диктора.", "bullet": True, "size": 16},
    {"text": "Высокая скорость: RTF ≈ 0,0078 на TPU; очищено порядка 1 млн ч аудио.", "bullet": True, "size": 16},
], top=4.25)

# =====================================================================
# Slide 5 — Базовая модель: недостатки
# =====================================================================
s = add_slide(); title(s, "Базовая модель: недостатки")
body(s, [
    {"text": "Закрытый энкодер USM (2 млрд параметров): ни веса, ни обученный очиститель признаков не опубликованы.", "bullet": True, "size": 18},
    {"text": "Прямое использование и воспроизведение Miipher-2 в открытых проектах невозможно.", "bullet": True, "size": 18},
    {"text": "Miipher-1 дополнительно требует текстовой расшифровки и идентификатора диктора на входе.", "bullet": True, "size": 18},
    {"text": "Реализация ориентирована на ускорители TPU — высокий порог входа.", "bullet": True, "size": 18},
    {"text": "Вывод: необходима открытая замена энкодера и подбор открытого вокодера, что и является задачей работы.", "bold": True, "size": 18, "space_before": 10},
], top=1.6, gap=12)

# =====================================================================
# Slide 6 — Формальная постановка
# =====================================================================
s = add_slide(); title(s, "Формальная постановка и обозначения")
body(s, [
    {"text": "Дан искажённый сигнал y (шум, реверберация, кодек). Модель искажения:", "size": 16}
], top=1.35, height=0.6)
picture(s, "m_degradation", left=2.0, top=1.95, max_w=9.3, max_h=0.95)
body(s, [{"text": "где h — импульсная характеристика помещения, n — шум, C — искажения тракта и кодека; уровень шума:", "size": 16}], top=3.0, height=0.6)
picture(s, "m_snr", left=3.0, top=3.55, max_w=7.3, max_h=1.1)
body(s, [{"text": "Требуется восстановить оценку чистого сигнала. Обозначения признаков:", "size": 16}], top=4.85, height=0.5)
picture(s, "m_notation", left=1.2, top=5.4, max_w=10.9, max_h=1.0)

# =====================================================================
# Slide 7 — Критерии качества
# =====================================================================
s = add_slide(); title(s, "Критерии качества решения")
body(s, [
    {"text": "PESQ — перцептивное качество (от −0,5 до 4,5), интрузивная метрика.", "bullet": True, "size": 17},
    {"text": "STOI — разборчивость речи (от 0 до 1), интрузивная.", "bullet": True, "size": 17},
    {"text": "SI-SDR — масштабно-инвариантное отношение сигнал/искажение, дБ:", "bullet": True, "size": 17},
], top=1.45, height=2.4)
picture(s, "m_sisdr", left=2.6, top=3.0, max_w=8.1, max_h=1.15)
body(s, [
    {"text": "DNSMOS — неинтрузивная нейросетевая оценка качества (без эталона).", "bullet": True, "size": 17},
    {"text": "RTF — коэффициент реального времени (быстродействие).", "bullet": True, "size": 17},
    {"text": "Требования: прирост PESQ ≥ 0,6 и STOI ≥ 0,05; RTF < 1; обучаемых параметров ≤ ~15 млн.", "bold": True, "size": 17, "space_before": 8},
], top=4.3, gap=8)

# =====================================================================
# Slide 8 — Предлагаемый метод: архитектура
# =====================================================================
s = add_slide(); title(s, "Предлагаемый метод: архитектура")
picture(s, "fig_arch", left=0.6, top=1.45, max_w=12.1, max_h=4.0,
        caption="Рис. 2. Параллельные адаптеры в замороженном энкодере; слой ℓ → вокодер")
body(s, [
    {"text": "Замороженный открытый SSL-энкодер + обучаемые параллельные адаптеры (очиститель) + открытый вокодер. Признаки выбранного слоя ℓ подаются на вокодер.", "size": 16},
], top=5.95)

# =====================================================================
# Slide 9 — Очиститель признаков
# =====================================================================
s = add_slide(); title(s, "Предлагаемый метод: очиститель признаков")
body(s, [{"text": "Очиститель — параллельные адаптеры с остаточной связью, встроенные в слои энкодера:", "size": 17}], top=1.4, height=0.7)
picture(s, "m_adapter", left=2.6, top=2.05, max_w=8.1, max_h=1.0)
body(s, [
    {"text": "Скрытая размерность узкого места 256, активация GELU; добавляет ~4,9 млн обучаемых параметров.", "bullet": True, "size": 16},
    {"text": "Энкодер заморожен — обучаются только адаптеры и вокодер.", "bullet": True, "size": 16},
    {"text": "Функция потерь очистителя (рассогласование признаков, λ = 1):", "size": 16, "space_before": 6},
], top=3.2, height=2.0)
picture(s, "m_fcloss", left=1.6, top=5.05, max_w=10.1, max_h=1.1)
body(s, [{"text": "Целевые признаки z* получаются прогоном чистого сигнала через тот же замороженный энкодер.", "size": 15, "italic": True, "color": GRAY}], top=6.3, height=0.6)

# =====================================================================
# Slide 10 — Выбор энкодеров
# =====================================================================
s = add_slide(); title(s, "Предлагаемый метод: выбор энкодеров")
body(s, [{"text": "Замена закрытого USM двумя открытыми SSL-энкодерами с близкой структурой:", "size": 17}], top=1.4, height=0.6)
table(s, [
    ["Энкодер", "Параметры", "Слоёв", "Данные предобуч.", "Особенность"],
    ["WavLM Base+", "95 млн", "12", "94 тыс. ч", "денойзинг при предобучении"],
    ["HuBERT Base", "95 млн", "12", "960 ч", "предсказание кластеров"],
    ["USM (закрыт)", "2 млрд", "32", "12 млн ч", "BEST-RQ, 300+ языков"],
], left=0.9, top=2.15, width=11.5, col_w=[2.2,1.4,1.0,2.0,3.2], size=15)
body(s, [
    {"text": "Гипотеза: благодаря моделированию шума при предобучении представления WavLM устойчивее к деградациям, чем у HuBERT.", "bold": True, "size": 17, "space_before": 4},
    {"text": "Оба энкодера: скрытая размерность 768, частота кадров 50 Гц (шаг 20 мс); используются в замороженном виде.", "size": 16},
], top=4.45, gap=10)

# =====================================================================
# Slide 11 — Исследование слоя
# =====================================================================
s = add_slide(); title(s, "Исследование слоя извлечения признаков")
picture(s, "fig_layers", left=1.6, top=1.45, max_w=9.0, max_h=4.0,
        caption="Рис. 3. Зависимость качества (PESQ) от номера слоя извлечения признаков")
body(s, [
    {"text": "Качество немонотонно зависит от слоя; максимум — на средних слоях.", "bullet": True, "size": 17},
    {"text": "Оптимум: ℓ = 7 (WavLM), ℓ = 6 (HuBERT).", "bullet": True, "size": 17},
    {"text": "Согласуется с SUPERB и с 13-м слоем USM (средняя треть сети).", "bullet": True, "size": 17},
], top=5.85)

# =====================================================================
# Slide 12 — Вокодер и конфигурации
# =====================================================================
s = add_slide(); title(s, "Предлагаемый метод: вокодер и конфигурации")
body(s, [{"text": "Вокодер обучается в генеративно-состязательном режиме (λ_fm = 2, λ_mel = 45):", "size": 16}], top=1.35, height=0.6)
picture(s, "m_vloss", left=2.4, top=1.95, max_w=8.5, max_h=1.0)
body(s, [
    {"text": "HiFi-GAN — однопроходный, быстрый; WaveFit — итеративный (5 итераций), естественнее.", "bullet": True, "size": 16},
    {"text": "Два энкодера × два вокодера = четыре исследуемые конфигурации:", "size": 16, "space_before": 4},
], top=3.1, height=1.1)
table(s, [
    ["Обозн.", "Энкодер", "Вокодер", "Комментарий"],
    ["K1", "HuBERT Base", "HiFi-GAN", "базовая быстрая конфигурация"],
    ["K2", "HuBERT Base", "WaveFit", "итеративный вокодер"],
    ["K3", "WavLM Base+", "HiFi-GAN", "устойчивый энкодер + быстрый вокодер"],
    ["K4", "WavLM Base+", "WaveFit", "наиболее близка к Miipher-2"],
], left=1.1, top=4.35, width=11.1, col_w=[1.0,2.0,1.6,4.0], size=14)

# =====================================================================
# Slide 13 — Параметры эксперимента
# =====================================================================
s = add_slide(); title(s, "Параметры вычислительного эксперимента")
body(s, [{"text": "Данные и гиперобучение", "bold": True, "size": 17}], top=1.35, height=0.4)
table(s, [
    ["Выборка", "Источник", "Объём", "Назначение"],
    ["Обучающая", "LibriTTS train", "240 ч", "обучение адаптеров и вокодера"],
    ["Валидационная", "LibriTTS dev", "12 ч", "подбор слоя и гиперпараметров"],
    ["Тест A", "VCTK", "8 ч", "невиданные дикторы"],
    ["Тест B", "LibriTTS test-clean", "10 ч", "сравнение результатов"],
], left=0.7, top=1.85, width=7.0, col_w=[1.6,2.2,0.9,2.6], size=12, height_row=0.5)
table(s, [
    ["Гиперпараметр", "Значение"],
    ["Оптимизатор", "AdamW"],
    ["Learning rate", "2·10⁻⁴"],
    ["Batch size", "16"],
    ["Сегмент / fs", "2 с / 16 кГц"],
    ["Шаги (адапт./вок.)", "200к / 400к"],
    ["GPU", "RTX 3090, 24 ГБ"],
], left=8.1, top=1.85, width=4.5, col_w=[2.5,2.0], size=12, height_row=0.5)
body(s, [
    {"text": "Деградации: шум (DEMAND, DNS Challenge), реверберация (RIR), сжатие кодеками; SNR ∈ [0; 20] дБ. Смесь формируется по правилу:", "size": 15},
], top=5.55, height=0.9)
picture(s, "m_mix", left=2.6, top=6.25, max_w=8.1, max_h=1.0)

# =====================================================================
# Slide 14 — Результаты: сравнение качества
# =====================================================================
s = add_slide(); title(s, "Результаты: сравнение качества")
table(s, [
    ["Конфигурация", "PESQ ↑", "STOI ↑", "SI-SDR, дБ ↑", "DNSMOS ↑"],
    ["Вход (шум)", "1,82", "0,841", "4,7", "2,41"],
    ["K1 (HuBERT+HiFi-GAN)", "2,71", "0,908", "13,9", "3,18"],
    ["K2 (HuBERT+WaveFit)", "2,78", "0,912", "14,3", "3,29"],
    ["K3 (WavLM+HiFi-GAN)", "**2,94**", "**0,927**", "15,2", "3,41"],
    ["K4 (WavLM+WaveFit)", "2,97", "0,929", "**15,6**", "**3,52**"],
], left=0.7, top=1.5, width=7.2, col_w=[3.0,1.3,1.3,1.7,1.5], size=12.5, height_row=0.55)
picture(s, "fig_tradeoff", left=8.0, top=1.5, max_w=4.9, max_h=3.6,
        caption="Рис. 4. Качество — быстродействие")
body(s, [
    {"text": "Все конфигурации дают прирост PESQ на 0,9–1,15 и STOI на 0,07–0,09 — выше целевых ориентиров.", "bullet": True, "size": 16},
    {"text": "Замена HuBERT → WavLM устойчиво улучшает все метрики; тип вокодера влияет на качество слабее.", "bullet": True, "size": 16},
], top=5.55)

# =====================================================================
# Slide 15 — Сравнение с базовыми, анализ
# =====================================================================
s = add_slide(); title(s, "Сравнение с базовыми методами и анализ")
table(s, [
    ["Метод", "PESQ ↑", "STOI ↑", "SI-SDR ↑", "DNSMOS ↑"],
    ["Вход (шум)", "1,82", "0,841", "4,7", "2,41"],
    ["Винеровская фильтрация", "2,11", "0,868", "9,8", "2,67"],
    ["DEMUCS", "2,64", "0,915", "**16,1**", "3,12"],
    ["K3 (наш подход)", "**2,94**", "**0,927**", "15,2", "**3,41**"],
], left=0.7, top=1.5, width=8.0, col_w=[3.2,1.2,1.2,1.4,1.4], size=13, height_row=0.55)
body(s, [
    {"text": "Анализ и выводы о гипотезе", "bold": True, "size": 17, "space_before": 4},
    {"text": "Классический фильтр не восстанавливает структуру сигнала; DEMUCS выше по SI-SDR, но уступает по перцептивным PESQ/DNSMOS (нет ресинтеза).", "bullet": True, "size": 15},
    {"text": "Конфигурации с HiFi-GAN ≈ в 4,4 раза быстрее аналогов с WaveFit при близком качестве.", "bullet": True, "size": 15},
    {"text": "Гипотеза подтверждена: открытые SSL-энкодеры заменяют USM; решающую роль играют выбор энкодера и слоя, а не вокодера.", "bullet": True, "bold": True, "size": 15},
], top=4.55, gap=7)

# =====================================================================
# Slide 16 — Заключение / Выносится на защиту
# =====================================================================
s = add_slide(); title(s, "Заключение")
body(s, [
    {"text": "Основной результат", "bold": True, "size": 17, "space_before": 0},
    {"text": "Открытая система восстановления речи на базе Miipher-2. Лучшая конфигурация — WavLM Base+ + HiFi-GAN (слой 7): PESQ 2,94; STOI 0,927; SI-SDR 15,2 дБ; DNSMOS 3,41; RTF 0,034.", "bullet": True, "size": 15},
    {"text": "Поставленная задача решена", "bold": True, "size": 17, "space_before": 6},
    {"text": "Все целевые метрики качества и быстродействия превышены; найдена конфигурация с наилучшим балансом.", "bullet": True, "size": 15},
    {"text": "Научная новизна", "bold": True, "size": 17, "space_before": 6},
    {"text": "Систематически исследовано влияние открытого SSL-энкодера и слоя извлечения признаков на качество восстановления в схеме Miipher-2.", "bullet": True, "size": 15},
    {"text": "Практическая значимость", "bold": True, "size": 17, "space_before": 6},
    {"text": "Автоматическое создание обучающих датасетов студийного качества из полностью открытых компонентов.", "bullet": True, "size": 15},
    {"text": "Планы на следующий этап", "bold": True, "size": 17, "space_before": 6},
    {"text": "WavLM Large и анализ его слоёв; сквозное (end-to-end) дообучение; многоязычные данные; субъективная оценка MOS.", "bullet": True, "size": 15},
], top=1.35, gap=3)

# =====================================================================
# Slide 18 — Список источников
# =====================================================================
s = add_slide(); title(s, "Список использованных источников")
refs = [
    "Koizumi Y. et al. Miipher-2: A Universal Speech Restoration Model for Million-Hour Scale Data Restoration // arXiv:2505.04457. — 2025.",
    "Koizumi Y. et al. Miipher: A Robust Speech Restoration Model Integrating Self-Supervised Speech and Text Representations // WASPAA. — 2023.",
    "Hsu W.-N. et al. HuBERT: Self-Supervised Speech Representation Learning by Masked Prediction of Hidden Units // IEEE/ACM TASLP. — 2021.",
    "Chen S. et al. WavLM: Large-Scale Self-Supervised Pre-Training for Full Stack Speech Processing // IEEE JSTSP. — 2022.",
    "Kong J., Kim J., Bae J. HiFi-GAN: GANs for Efficient and High Fidelity Speech Synthesis // NeurIPS. — 2020.",
    "Koizumi Y. et al. WaveFit: An Iterative and Non-Autoregressive Neural Vocoder // IEEE SLT. — 2023.",
    "Baevski A. et al. wav2vec 2.0: A Framework for Self-Supervised Learning of Speech Representations // NeurIPS. — 2020.",
    "Zhang Y. et al. Google USM: Scaling Automatic Speech Recognition Beyond 100 Languages // arXiv:2303.01037. — 2023.",
]
body(s, [{"text": f"{i+1}. {r}", "italic": True, "size": 14.5, "space_after": 7} for i, r in enumerate(refs)],
     top=1.5, gap=7)

prs.save("VKR_Vambrikov_presentation.pptx")
print("saved VKR_Vambrikov_presentation.pptx; slides:", len(prs.slides))
